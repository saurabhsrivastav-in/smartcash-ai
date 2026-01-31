import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from datetime import datetime, timedelta
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO

# --- FIXED REPORT GENERATION ENGINES ---

def generate_pdf(df, mode_name, liquidity):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", 'B', 16)
    pdf.cell(190, 10, "SmartCash AI: Executive Treasury Report", ln=True, align='C')
    
    pdf.set_font("Helvetica", size=10)
    pdf.cell(190, 10, f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}", ln=True, align='C')
    pdf.ln(10)
    
    pdf.set_font("Helvetica", 'B', 12)
    pdf.cell(95, 10, f"Scenario: {mode_name}")
    pdf.cell(95, 10, f"Net Liquidity: ${liquidity/1e6:.2f}M", ln=True)
    pdf.ln(5)
    
    # Table Header
    pdf.set_fill_color(22, 27, 34)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(40, 10, "Invoice", 1, 0, 'C', True)
    pdf.cell(60, 10, "Customer", 1, 0, 'C', True)
    pdf.cell(45, 10, "Amount", 1, 0, 'C', True)
    pdf.cell(45, 10, "Due Date", 1, 1, 'C', True)
    
    # Table Body
    pdf.set_text_color(0, 0, 0)
    pdf.set_font("Helvetica", size=10)
    for _, row in df.head(20).iterrows():
        pdf.cell(40, 10, str(row['Invoice_ID']), 1)
        pdf.cell(60, 10, str(row['Customer'])[:25], 1)
        pdf.cell(45, 10, f"{row['Amount_Remaining']:,.2f}", 1)
        pdf.cell(45, 10, str(row['Due_Date']), 1, 1)
        
    # --- INDENTATION FIXED BELOW ---
    pdf_output = pdf.output()
    if isinstance(pdf_output, (bytearray, str)):
        if isinstance(pdf_output, str):
            return pdf_output.encode('latin-1')
        return bytes(pdf_output)
    return pdf_output

def generate_pptx(df, mode_name, liquidity):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    
    title = slide.shapes.title
    title.text = f"SmartCash AI Executive Summary"
    
    # KPI Stats
    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content.text_frame
    p = tf.add_paragraph()
    p.text = f"Scenario Mode: {mode_name}"
    p.font.size = Pt(24)
    
    p2 = tf.add_paragraph()
    p2.text = f"Risk-Adjusted Liquidity: ${liquidity/1e6:.2f}M"
    p2.font.bold = True
    p2.font.size = Pt(32)
    
    # CRITICAL FIX: Use BytesIO correctly
    binary_output = BytesIO()
    prs.save(binary_output)
    binary_output.seek(0) # Move to start of buffer
    return binary_output.getvalue()

# --- 1. STABILITY INITIALIZATION ---
if 'audit' not in st.session_state:
    st.session_state.audit = []
if 'ledger' not in st.session_state:
    customers = ['Tesla', 'EcoEnergy', 'GlobalBlue', 'TechRetail', 'Quantum Dyn', 'Alpha Log', 'Nordic Oil', 'Sino Tech', 'Indo Power', 'Euro Mart']
    entities = ['1000 (US)', '2000 (EU)', '3000 (UK)']
    currencies = {'1000 (US)': 'USD', '2000 (EU)': 'EUR', '3000 (UK)': 'GBP'}
    ratings = ['AAA', 'AA', 'A', 'B', 'C', 'D']
    
    inv_data = []
    for i in range(300):
        ent = np.random.choice(entities)
        # FIX 1: Increase the range to 10M - 35M
        amt = np.random.uniform(10000000, 35000000) 
        
        due = datetime(2026, 1, 30) - timedelta(days=np.random.randint(-30, 90))
        inv_data.append({
            'Invoice_ID': f"INV-{1000+i}",
            'Company_Code': ent,
            'Customer': np.random.choice(customers),
            'Amount_Remaining': round(amt, 2),
            'Currency': currencies[ent],
            'ESG_Score': np.random.choice(ratings),
            'Due_Date': due.strftime('%Y-%m-%d'),
            'Status': 'Overdue' if due < datetime(2026, 1, 30) else 'Open',
            'Is_Disputed': False
        })
    st.session_state.ledger = pd.DataFrame(inv_data)

# --- 2. EXECUTIVE THEME ---
st.set_page_config(page_title="SmartCash AI | C-Suite", page_icon="🏛️", layout="wide")
st.markdown("""
    <style>
    .stApp { background-color: #0d1117; color: #c9d1d9; }
    [data-testid="stMetricValue"] { font-size: 28px; color: #58a6ff; font-weight: 800; }
    .stMetric { background: #161b22; border: 1px solid #30363d; border-radius: 12px; padding: 10px; }
    </style>
    """, unsafe_allow_html=True)

# --- 3. SIDEBAR & MACRO TOGGLES ---
with st.sidebar:
    st.title("🛡️ Risk Controls")
    # Interactive Toggles that affect metrics globally
    bad_debt_provision = st.slider("Bad Debt Provision (%)", 0, 20, 5)
    risk_weighting = st.toggle("Enable Risk-Weighted Valuation", value=True)
    st.divider()
    st.info(f"System Time: {datetime.now().strftime('%Y-%m-%d %H:%M')}")

# --- 4. TOP NAVIGATION ---
st.title("🏛️ Executive Treasury Intelligence")
suggestion_list = ["Consolidated"] + sorted(st.session_state.ledger['Customer'].unique().tolist())

h_col1, h_col2, h_col3 = st.columns([3, 2, 1])
with h_col1:
    search_selection = st.selectbox("🎯 Strategic Entity Search", options=suggestion_list, index=0)
with h_col2:
    st.write(" ") 
    mode = st.segmented_control("View Mode", ["Actuals", "AI Forecast", "Stress Test"], default="Actuals")
with h_col3:
    st.write(" ")
    if st.button("🔄 Refresh Data"): st.rerun()

st.divider()

# --- 5. DATA FILTERING & SCENARIO ENGINE ---
view_df = st.session_state.ledger.copy()

# Filter by Customer
if search_selection != "Consolidated":
    view_df = view_df[view_df['Customer'] == search_selection]

# View Mode Logic (THE FIX: Defining weights for ALL paths)
if mode == "AI Forecast":
    # 1. AI Logic: Predict 2% haircut on all values
    view_df['Amount_Remaining'] = view_df['Amount_Remaining'] * 0.98 
    # 2. AI Weights: Slightly more optimistic than 'Actuals'
    weights = {'AAA':0.99, 'AA':0.97, 'A':0.92, 'B':0.85, 'C':0.70, 'D':0.50}
    st.sidebar.warning("🤖 AI Mode: Adjusting for predicted defaults.")

elif mode == "Stress Test":
    # 1. Stress Logic: Aggressive risk weighting for downturn
    weights = {'AAA':0.90, 'AA':0.80, 'A':0.70, 'B':0.40, 'C':0.20, 'D':0.05}
    st.sidebar.error("🔥 Stress Test: 20% Market Downturn applied.")

else:
    # 1. Actuals Logic: Standard institutional weights
    weights = {'AAA':0.98, 'AA':0.95, 'A':0.90, 'B':0.80, 'C':0.60, 'D':0.40}

# --- 6. C-SUITE METRICS (Dynamic Calculations) ---
total_val = view_df['Amount_Remaining'].sum()

if mode == "Stress Test" or risk_weighting:
    # Calculate liquidity based on scenario-specific weights
    net_collectible = (view_df['Amount_Remaining'] * view_df['ESG_Score'].map(weights)).sum()
else:
    net_collectible = total_val * (1 - (bad_debt_provision/100))

# Update Metrics
m1, m2, m3, m4 = st.columns(4)
m1.metric("Risk-Adjusted Net Liquidity", f"${(net_collectible/1e6):.2f}M", 
          delta=f"{(net_collectible/total_val*100):.1f}% Realizable")
m2.metric("Working Capital Pool", f"${(total_val/1e6):.1f}M")
m3.metric("Weighted DSO", "34 Days" if mode == "AI Forecast" else "38 Days", 
          delta="-4 Days" if mode == "AI Forecast" else "Flat")
m4.metric("Capital at Risk", f"${((total_val - net_collectible)/1e6):.2f}M", 
          delta="Increased" if mode == "Stress Test" else "Stable", delta_color="inverse")

# --- 7. GRAPHICAL INTELLIGENCE ---
tab_charts, tab_velocity, tab_entity, tab_stress = st.tabs(["📊 Exposure Analytics", "⏳ Liquidity Timeline", "🏢 Entity Exposure", "🔥 Stress Matrix"])

with tab_charts:
    c1, c2 = st.columns([2, 1])
    with c1:
        st.subheader("🛡️ Strategic Risk Radar")
        fig_s = px.sunburst(view_df, path=['Company_Code', 'ESG_Score', 'Customer'], values='Amount_Remaining', color='ESG_Score', color_discrete_map={'AAA':'#238636', 'AA':'#2ea043', 'A':'#d29922', 'B':'#db6d28', 'C':'#f85149', 'D':'#b62323'})
        fig_s.update_layout(template="plotly_dark", height=500, margin=dict(t=10, b=10, l=10, r=10))
        st.plotly_chart(fig_s, use_container_width=True)
    with c2:
        st.subheader("⏳ Institutional Ageing")
        def get_bucket(d):
            days = (datetime(2026, 1, 30) - datetime.strptime(d, '%Y-%m-%d')).days
            return "0-30" if days <= 30 else "31-60" if days <= 60 else "61-90" if days <= 90 else "90+"
        ov = view_df[view_df['Status'] == 'Overdue'].copy()
        if not ov.empty:
            ov['Bucket'] = ov['Due_Date'].apply(get_bucket)
            age_data = ov.groupby('Bucket')['Amount_Remaining'].sum().reset_index()
            fig_age = px.bar(age_data, x='Bucket', y='Amount_Remaining', color='Bucket', color_discrete_sequence=px.colors.sequential.Reds_r)
            fig_age.update_layout(template="plotly_dark", height=500, showlegend=False)
            st.plotly_chart(fig_age, use_container_width=True)

with tab_velocity:
    st.subheader("📈 Projected Cash Inflow Velocity")
    view_df['Due_Date_DT'] = pd.to_datetime(view_df['Due_Date'])
    velocity_data = view_df.groupby(pd.Grouper(key='Due_Date_DT', freq='W'))['Amount_Remaining'].sum().reset_index()
    fig_v = px.area(velocity_data, x='Due_Date_DT', y='Amount_Remaining', color_discrete_sequence=['#58a6ff'])
    fig_v.update_layout(template="plotly_dark", xaxis_title="Projection Horizon", yaxis_title="Volume ($)")
    st.plotly_chart(fig_v, use_container_width=True)

with tab_entity:
    st.subheader("🏢 Strategic Entity Risk Profiling")
    st.markdown("---")
    
    risk_colors = {'AAA':'#238636', 'AA':'#2ea043', 'A':'#d29922', 'B':'#db6d28', 'C':'#f85149', 'D':'#b62323'}
    entity_stats = view_df.copy()
    entity_stats['Due_DT'] = pd.to_datetime(entity_stats['Due_Date'])
    entity_stats['Days_Late'] = (pd.to_datetime('2026-01-30') - entity_stats['Due_DT']).dt.days.clip(lower=0)
    
    entity_analysis = entity_stats.groupby(['Customer', 'ESG_Score', 'Company_Code']).agg({
        'Amount_Remaining': 'sum',
        'Invoice_ID': 'count',
        'Days_Late': 'mean'
    }).reset_index()

    # --- ENLARGED LAYOUT: 3:1 Ratio ---
    col_matrix, col_cards = st.columns([3, 1])

    with col_matrix:
        st.write("#### 🎯 Strategic Risk Matrix")
        fig_bubble = px.scatter(
            entity_analysis,
            x="Days_Late",
            y="Amount_Remaining",
            size="Amount_Remaining",
            color="ESG_Score",
            hover_name="Customer",
            color_discrete_map=risk_colors,
            labels={"Days_Late": "Avg Days Overdue", "Amount_Remaining": "Total Exposure ($)"},
            template="plotly_dark",
            size_max=60 
        )
        
        fig_bubble.update_layout(
            height=650, 
            margin=dict(l=20, r=20, t=40, b=20),
            xaxis=dict(gridcolor="#30363d"),
            yaxis=dict(gridcolor="#30363d")
        )
        
        fig_bubble.add_vline(x=30, line_dash="dash", line_color="#f85149", annotation_text="30D Threshold")
        fig_bubble.add_hrect(y0=total_val*0.1, y1=entity_analysis['Amount_Remaining'].max()*1.2, 
                             fillcolor="red", opacity=0.05, annotation_text="CONCENTRATION LIMIT")
        
        st.plotly_chart(fig_bubble, use_container_width=True, key="entity_bubble_enlarged")

    with col_cards:
        st.write("#### 🛡️ Priority Watchlist")
        top_risks = entity_analysis.sort_values(by='Amount_Remaining', ascending=False).head(5)
        
        for _, row in top_risks.iterrows():
            border_color = risk_colors.get(row['ESG_Score'], '#58a6ff')
            st.markdown(f"""
            <div style="background:#161b22; padding:12px; border-radius:10px; border-left: 5px solid {border_color}; margin-bottom:10px; border-top: 1px solid #30363d;">
                <p style="margin:0; font-weight:bold; font-size:14px;">{row['Customer']}</p>
                <p style="margin:0; font-size:18px; color:#58a6ff; font-weight:800;">${row['Amount_Remaining']/1e6:.2f}M</p>
                <p style="margin:0; font-size:11px; color:#8b949e;">{int(row['Days_Late'])} Days Overdue</p>
            </div>
            """, unsafe_allow_html=True)

    st.divider()
    b_col1, b_col2 = st.columns(2)
    with b_col1:
        st.write("#### 🌍 Regional Allocation")
        fig_reg = px.bar(entity_analysis, x="Company_Code", y="Amount_Remaining", color="ESG_Score", color_discrete_map=risk_colors, template="plotly_dark")
        st.plotly_chart(fig_reg, use_container_width=True, key="reg_bar_vfinal")
    with b_col2:
        st.write("#### 📊 Portfolio Rating Mix")
        fig_pie = px.pie(entity_analysis, names="ESG_Score", values="Amount_Remaining", hole=0.5, color="ESG_Score", color_discrete_map=risk_colors, template="plotly_dark")
        st.plotly_chart(fig_pie, use_container_width=True, key="pie_final_vfinal")

with tab_stress:
    st.subheader("Strategic Liquidity Stress Matrix (FX Volatility vs. Hedging)")
    fx_range = np.array([-10, -5, 0, 5, 10])
    hedge_range = np.array([0, 25, 50, 75, 100])
    base_liq = net_collectible / 1e6
    z_data = [[round(base_liq * (1 + (fx/100) * (1 - (h/100))), 2) for h in hedge_range] for fx in fx_range]
    fig_h = go.Figure(data=go.Heatmap(z=z_data, x=[f"{h}% Hedge" for h in hedge_range], y=[f"{fx}% FX Vol" for fx in fx_range], colorscale='RdYlGn', text=z_data, texttemplate="$%{text}M"))
    fig_h.update_layout(template="plotly_dark", height=450)
    st.plotly_chart(fig_h, use_container_width=True, key="stress_heatmap")
    
    # --- BOTTOM ROW: Regional Diversification & Volume Distribution ---
    col_b1, col_b2 = st.columns([1, 1])
    
    with col_b1:
        st.write("#### 🌍 Regional/Company Code Diversification")
        fig_bar = px.bar(
            entity_analysis, 
            x="Company_Code", 
            y="Amount_Remaining", 
            color="ESG_Score",
            barmode="stack",
            template="plotly_dark",
            color_discrete_map=risk_colors
        )
        st.plotly_chart(fig_bar, use_container_width=True, key="entity_regional_bar_v3")
        
    with col_b2:
        st.write("#### 📊 Risk Distribution by Invoice Volume")
        fig_pie = px.pie(
            entity_analysis,
            names="ESG_Score",
            values="Invoice_ID",
            color="ESG_Score",
            hole=0.4,
            color_discrete_map=risk_colors,
            template="plotly_dark"
        )
        fig_pie.update_traces(textposition='inside', textinfo='percent+label')
        st.plotly_chart(fig_pie, use_container_width=True, key="entity_volume_pie")
    
    # Geographical/Company Code Diversification
    st.write("#### 🌍 Regional/Company Code Diversification")
    fig_bar = px.bar(
        entity_analysis, 
        x="Company_Code", 
        y="Amount_Remaining", 
        color="ESG_Score",
        barmode="stack",
        template="plotly_dark",
        color_discrete_map=risk_colors
    )
    st.plotly_chart(fig_bar, use_container_width=True, key="entity_regional_bar")
    
    # Sector/Code breakdown for geographical/entity risk
    st.write("#### 🌍 Regional/Company Code Diversification")
    fig_bar = px.bar(
        entity_analysis, 
        x="Company_Code", 
        y="Amount_Remaining", 
        color="ESG_Score",
        barmode="stack",
        template="plotly_dark",
        color_discrete_map={'AAA':'#238636', 'AA':'#2ea043', 'A':'#d29922', 'B':'#db6d28', 'C':'#f85149', 'D':'#b62323'}
    )
    st.plotly_chart(fig_bar, use_container_width=True)

with tab_stress:
    st.subheader("Strategic Liquidity Stress Matrix (FX Volatility vs. Hedging)")
    fx_range = np.array([-10, -5, 0, 5, 10])
    hedge_range = np.array([0, 25, 50, 75, 100])
    base_liq = net_collectible / 1e6
    z_data = [[round(base_liq * (1 + (fx/100) * (1 - (h/100))), 2) for h in hedge_range] for fx in fx_range]
    fig_h = go.Figure(data=go.Heatmap(z=z_data, x=[f"{h}% Hedge" for h in hedge_range], y=[f"{fx}% FX Vol" for fx in fx_range], colorscale='RdYlGn', text=z_data, texttemplate="$%{text}M"))
    fig_h.update_layout(template="plotly_dark", height=450)
    st.plotly_chart(fig_h, use_container_width=True)

# --- 8. EXECUTIVE ACTION PANEL ---
st.subheader("⚡ Executive Action Thresholds")
cols = st.columns(3)
with cols[0]:
    if st.button("🚀 Authorize AR Discounting", use_container_width=True):
        st.toast("Discounting Facility Activated for High-Risk Buckets")
with cols[1]:
    if st.button("🛡️ Trigger FX Hedge", use_container_width=True):
        st.toast("Hedge Ratio increased to 75% for Active Portfolio")
with cols[2]:
    if st.button("📩 Escalated Board Report", use_container_width=True):
        st.toast("Executive Summary PDF generated and sent.")

# --- 9. UI DOWNLOAD SECTION ---

st.divider()
st.subheader("📤 Export Intelligence")
d_col1, d_col2 = st.columns(2)

with d_col1:
    # We use a unique key and ensure generate_pdf returns clean bytes
    try:
        pdf_bytes = generate_pdf(view_df, mode, net_collectible)
        st.download_button(
            label="📥 Download Executive PDF",
            data=pdf_bytes,
            file_name=f"SmartCash_Report_{mode}.pdf",
            mime="application/pdf",
            use_container_width=True,
            key="pdf_download_btn"
        )
    except Exception as e:
        st.error("PDF Engine encountered a buffer error. Please refresh.")

with d_col2:
    try:
        pptx_bytes = generate_pptx(view_df, mode, net_collectible)
        st.download_button(
            label="📊 Download Board PPTX",
            data=pptx_bytes,
            file_name=f"SmartCash_Deck_{mode}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True,
            key="pptx_download_btn"
        )
    except Exception as e:
        st.error("PPTX Engine encountered a buffer error.")
